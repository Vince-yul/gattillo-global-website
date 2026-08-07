#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from html.parser import HTMLParser

# Color scheme
DARK_BG = RGBColor(14, 12, 10)
LIGHT_BG = RGBColor(245, 243, 240)
ACCENT = RGBColor(212, 166, 68)
BONE = RGBColor(237, 230, 216)
DARK_TEXT = RGBColor(26, 24, 21)

# Slide content structure
slides_data = [
    {
        "title": "The Innovation\nOpportunity",
        "subtitle": "A conversation with Chinese arcade manufacturers about leading the industry, not following it.",
        "meta": "GATTILLO GLOBAL",
        "bg": "dark",
        "type": "title"
    },
    {
        "title": "This exists in China today",
        "subtitle": "The budget segment is alive and thriving. Commoditized products, racing to the bottom on price.",
        "columns": [
            {
                "heading": "Budget Arcade",
                "items": ["Generic machines", "Low-cost design", "Indistinguishable products", "Race to the bottom on price"]
            },
            {
                "heading": "Budget Cars",
                "items": ["$3,000–$7,000 vehicles", "Basic mobility, nothing more", "Same commoditized playbook", "Manufacturers stuck competing on cost"]
            }
        ],
        "bg": "dark",
        "type": "comparison"
    },
    {
        "title": "Today in China, companies are building premium products at scale.",
        "boxes": [
            {"name": "BYD", "desc": "Built for premium. Superior design, technology, and scale. Competing globally against Tesla."},
            {"name": "NIO", "desc": "Luxury EV brand. Premium experience. Proving Chinese innovation can lead, not follow."},
            {"name": "And More", "desc": "New entrants building in the premium tier. It's happening now. In China."}
        ],
        "bg": "dark",
        "type": "three_col"
    },
    {
        "title": "The EV Revolution Is Our Blueprint",
        "sections": [
            {
                "heading": "2015: The Budget Segment",
                "text": "Generic. Commodity. Undifferentiated. This is arcade games today."
            },
            {
                "heading": "2024: Innovation Leaped Ahead",
                "text": "Better design. Better performance. Lower cost. This is where we go together."
            }
        ],
        "bg": "light",
        "type": "two_section"
    },
    {
        "title": "What Gattillo Global Offers",
        "items": [
            "Novel game concepts (not rehashed classics)",
            "Player psychology & engagement design",
            "Access to international IPs",
            "Operator profitability frameworks",
            "Market positioning & brand architecture"
        ],
        "bg": "dark",
        "type": "list"
    },
    {
        "title": "Your Contribution",
        "items": [
            "Manufacturing excellence at scale",
            "Cost efficiency (10–40% below Western production)",
            "Supply chain agility & rapid iteration",
            "Regional market knowledge & operator relationships",
            "Proven ability to deliver at volume"
        ],
        "bg": "dark",
        "type": "list"
    },
    {
        "title": "This Is A Partnership",
        "partnerships": [
            {
                "num": "1",
                "heading": "Co-Design",
                "desc": "You're not a vendor. We develop together. Your manufacturing insight shapes the final product."
            },
            {
                "num": "2",
                "heading": "Shared Risk",
                "desc": "We invest in concept validation and business development, gaining market insights before full production. You invest in tooling and capacity. Aligned incentives."
            },
            {
                "num": "3",
                "heading": "Market Ownership",
                "desc": "You lead in Asia, build your brand. Gattillo captures global and premium markets. Growing pie, not zero-sum."
            }
        ],
        "bg": "dark",
        "type": "partnership"
    },
    {
        "title": "Very soon",
        "points": [
            ("Arcade games designed in North America and built in China.", "Not the cheap version of a Western design, but the better version at a fraction of the cost."),
            ("We have shifted the entire market.", "SEGA and Namco are no longer the default. We are competing on design and novelty, not racing to the bottom on price."),
            ("That's the BYD moment for arcade games.", "And it starts with this conversation.")
        ],
        "bg": "light",
        "type": "vision"
    },
    {
        "title": "Here's What We're Asking",
        "cta_label": "The Conversation",
        "cta_heading": "Are you interested in being the Chinese innovator in arcade games?",
        "cta_text": "Not just the low-cost manufacturer. The brand that changes what's possible.\n\nIf yes, let's explore the first co-designed game together.",
        "bg": "dark",
        "type": "cta"
    },
    {
        "title": "Let's Talk",
        "subtitle": "We'll bring the game concepts. You bring the manufacturing excellence. Together, we'll prove that China can lead the industry, not just fill it.",
        "footer": "GATTILLO GLOBAL\nReaching out to shape the future of play.",
        "bg": "dark",
        "type": "title"
    }
]

def add_title_slide(prs, slide_data):
    """Add a title/closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG if slide_data["bg"] == "dark" else LIGHT_BG

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = BONE if slide_data["bg"] == "dark" else DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    if "subtitle" in slide_data and slide_data["subtitle"]:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = slide_data["subtitle"]
        p.font.size = Pt(24)
        p.font.color.rgb = BONE if slide_data["bg"] == "dark" else DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

    # Add footer
    if "footer" in slide_data and slide_data["footer"]:
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1))
        footer_frame = footer_box.text_frame
        footer_frame.word_wrap = True
        p = footer_frame.paragraphs[0]
        p.text = slide_data["footer"]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER

    # Add meta if exists
    if "meta" in slide_data and slide_data["meta"]:
        meta_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        meta_frame = meta_box.text_frame
        p = meta_frame.paragraphs[0]
        p.text = slide_data["meta"]
        p.font.size = Pt(12)
        p.font.color.rgb = BONE if slide_data["bg"] == "dark" else DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

def add_comparison_slide(prs, slide_data):
    """Add a comparison slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BONE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = slide_data["subtitle"]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Two columns
    left_x = Inches(0.75)
    right_x = Inches(5.25)
    col_width = Inches(4)

    for idx, col in enumerate(slide_data["columns"]):
        x = left_x if idx == 0 else right_x

        # Column heading
        heading_box = slide.shapes.add_textbox(x, Inches(2.5), col_width, Inches(0.6))
        heading_frame = heading_box.text_frame
        p = heading_frame.paragraphs[0]
        p.text = col["heading"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        # Items
        items_box = slide.shapes.add_textbox(x, Inches(3.3), col_width, Inches(3.5))
        items_frame = items_box.text_frame
        items_frame.word_wrap = True

        for item_idx, item in enumerate(col["items"]):
            if item_idx == 0:
                p = items_frame.paragraphs[0]
            else:
                p = items_frame.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(18)
            p.font.color.rgb = BONE
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0

def add_list_slide(prs, slide_data):
    """Add a simple list slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BONE
    p.alignment = PP_ALIGN.CENTER

    # List
    list_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(5))
    list_frame = list_box.text_frame
    list_frame.word_wrap = True

    for idx, item in enumerate(slide_data["items"]):
        if idx == 0:
            p = list_frame.paragraphs[0]
        else:
            p = list_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = BONE
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0

def add_partnership_slide(prs, slide_data):
    """Add partnership slide with three columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BONE
    p.alignment = PP_ALIGN.CENTER

    # Three columns
    col_width = Inches(2.8)
    col_x_positions = [Inches(0.5), Inches(3.5), Inches(6.5)]

    for idx, partnership in enumerate(slide_data["partnerships"]):
        x = col_x_positions[idx]
        y_start = Inches(1.8)

        # Number
        num_box = slide.shapes.add_textbox(x, y_start, col_width, Inches(0.6))
        num_frame = num_box.text_frame
        p = num_frame.paragraphs[0]
        p.text = partnership["num"]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER

        # Heading
        heading_box = slide.shapes.add_textbox(x, y_start + Inches(0.7), col_width, Inches(0.6))
        heading_frame = heading_box.text_frame
        heading_frame.word_wrap = True
        p = heading_frame.paragraphs[0]
        p.text = partnership["heading"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BONE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(x, y_start + Inches(1.5), col_width, Inches(3.5))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = partnership["desc"]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

def add_vision_slide(prs, slide_data):
    """Add vision slide with callout points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Points
    y_pos = Inches(1.7)
    for heading, text in slide_data["points"]:
        # Heading
        heading_box = slide.shapes.add_textbox(Inches(1), y_pos, Inches(8), Inches(0.4))
        heading_frame = heading_box.text_frame
        heading_frame.word_wrap = True
        p = heading_frame.paragraphs[0]
        p.text = heading
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        # Text
        text_box = slide.shapes.add_textbox(Inches(1), y_pos + Inches(0.4), Inches(8), Inches(0.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT

        y_pos += Inches(1.3)

def add_cta_slide(prs, slide_data):
    """Add call-to-action slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BONE
    p.alignment = PP_ALIGN.CENTER

    # CTA Box
    box_left = Inches(1.75)
    box_top = Inches(2.2)
    box_width = Inches(6.5)
    box_height = Inches(3.5)

    # Draw box border
    shape = slide.shapes.add_shape(1, box_left, box_top, box_width, box_height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(40, 40, 40)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(2)

    # CTA Label
    label_box = slide.shapes.add_textbox(box_left + Inches(0.3), box_top + Inches(0.3), box_width - Inches(0.6), Inches(0.35))
    label_frame = label_box.text_frame
    p = label_frame.paragraphs[0]
    p.text = slide_data["cta_label"]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # CTA Heading
    heading_box = slide.shapes.add_textbox(box_left + Inches(0.3), box_top + Inches(0.8), box_width - Inches(0.6), Inches(1.2))
    heading_frame = heading_box.text_frame
    heading_frame.word_wrap = True
    p = heading_frame.paragraphs[0]
    p.text = slide_data["cta_heading"]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BONE

    # CTA Text
    text_box = slide.shapes.add_textbox(box_left + Inches(0.3), box_top + Inches(2.1), box_width - Inches(0.6), Inches(1.2))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = slide_data["cta_text"]
    p.font.size = Pt(14)
    p.font.color.rgb = BONE

def add_three_col_slide(prs, slide_data):
    """Add three column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BONE
    p.alignment = PP_ALIGN.CENTER

    # Three columns
    col_width = Inches(2.8)
    col_x_positions = [Inches(0.5), Inches(3.5), Inches(6.5)]

    for idx, box in enumerate(slide_data["boxes"]):
        x = col_x_positions[idx]
        y_start = Inches(1.8)

        # Name
        name_box = slide.shapes.add_textbox(x, y_start, col_width, Inches(0.6))
        name_frame = name_box.text_frame
        p = name_frame.paragraphs[0]
        p.text = box["name"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = ACCENT if box["name"] != "And More" else BONE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(x, y_start + Inches(1.2), col_width, Inches(3.5))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = box["desc"]
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

def add_two_section_slide(prs, slide_data):
    """Add two section comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = slide_data["title"]
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Two sections
    col_width = Inches(4)
    left_x = Inches(0.75)
    right_x = Inches(5.25)

    for idx, section in enumerate(slide_data["sections"]):
        x = left_x if idx == 0 else right_x
        y_start = Inches(1.8)

        # Heading
        heading_box = slide.shapes.add_textbox(x, y_start, col_width, Inches(0.8))
        heading_frame = heading_box.text_frame
        heading_frame.word_wrap = True
        p = heading_frame.paragraphs[0]
        p.text = section["heading"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = ACCENT if "2024" in section["heading"] else DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

        # Text
        text_box = slide.shapes.add_textbox(x, y_start + Inches(1.2), col_width, Inches(3.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = section["text"]
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add all slides
for slide_data in slides_data:
    slide_type = slide_data.get("type", "title")

    if slide_type == "title":
        add_title_slide(prs, slide_data)
    elif slide_type == "comparison":
        add_comparison_slide(prs, slide_data)
    elif slide_type == "list":
        add_list_slide(prs, slide_data)
    elif slide_type == "three_col":
        add_three_col_slide(prs, slide_data)
    elif slide_type == "two_section":
        add_two_section_slide(prs, slide_data)
    elif slide_type == "partnership":
        add_partnership_slide(prs, slide_data)
    elif slide_type == "vision":
        add_vision_slide(prs, slide_data)
    elif slide_type == "cta":
        add_cta_slide(prs, slide_data)

# Save presentation
output_path = "/Users/Vincent/Documents/Projects/Gattillo-Global-Website/gattillo-china-innovation-deck.pptx"
prs.save(output_path)
print(f"✓ PowerPoint presentation saved to: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
